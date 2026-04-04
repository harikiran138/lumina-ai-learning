from __future__ import annotations

import csv
import io
import json
import os
import re
import tempfile
from collections import Counter
from pathlib import Path
from statistics import median
from typing import Any, Dict, List, Optional

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pypdf import PdfReader

from app.core.logging import structlog
from app.services.ppt_generator import PPTGenerator
from app.services.storage import storage_service
from ai_engine.llm import get_llm_provider, is_provider_error

try:
    import pdfplumber
except ImportError:  # pragma: no cover
    pdfplumber = None

log = structlog.get_logger()


class UnitPDFParser:
    def parse(self, file_path: str, source_filename: str = "unit.pdf") -> Dict[str, Any]:
        parsed = self._parse_with_pdfplumber(file_path) if pdfplumber is not None else None
        if parsed:
            return parsed
        return self._parse_with_pypdf(file_path, source_filename)

    def _parse_with_pdfplumber(self, file_path: str) -> Optional[Dict[str, Any]]:
        try:
            with pdfplumber.open(file_path) as pdf:
                page_payloads: List[Dict[str, Any]] = []
                all_sizes: List[float] = []
                for page_index, page in enumerate(pdf.pages):
                    lines = self._extract_lines_from_page(page)
                    all_sizes.extend([line["size"] for line in lines if line.get("size")])
                    page_payloads.append(
                        {
                            "page_number": page_index + 1,
                            "lines": lines,
                            "tables": self._extract_tables(page, page_index + 1),
                            "images": self._extract_images(page, page_index + 1),
                        }
                    )
                heading_threshold = (median(all_sizes) + 1.5) if all_sizes else 14.0
                return self._assemble_structure(page_payloads, heading_threshold, Path(file_path).stem)
        except Exception as exc:
            log.warning("pdfplumber_parse_failed", file_path=file_path, error=str(exc))
            return None

    def _parse_with_pypdf(self, file_path: str, source_filename: str) -> Dict[str, Any]:
        reader = PdfReader(file_path)
        page_payloads: List[Dict[str, Any]] = []
        for page_index, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            lines = []
            for raw_line in text.splitlines():
                clean = raw_line.strip()
                if clean:
                    lines.append({"text": clean, "size": 12.0})
            page_payloads.append(
                {"page_number": page_index + 1, "lines": lines, "tables": [], "images": []}
            )
        return self._assemble_structure(page_payloads, 13.0, Path(source_filename).stem)

    def _extract_lines_from_page(self, page: Any) -> List[Dict[str, Any]]:
        words = page.extract_words(use_text_flow=True, extra_attrs=["size"]) or []
        grouped: Dict[int, List[Dict[str, Any]]] = {}
        for word in words:
            top_key = int(round(float(word.get("top", 0.0)) / 3.0) * 3)
            grouped.setdefault(top_key, []).append(word)

        lines: List[Dict[str, Any]] = []
        for top in sorted(grouped):
            row_words = sorted(grouped[top], key=lambda item: float(item.get("x0", 0.0)))
            text = " ".join(word.get("text", "").strip() for word in row_words).strip()
            if not text:
                continue
            size_values = [float(word.get("size", 12.0)) for word in row_words if word.get("size")]
            lines.append({"text": text, "size": sum(size_values) / len(size_values) if size_values else 12.0})
        return lines

    def _extract_tables(self, page: Any, page_number: int) -> List[Dict[str, Any]]:
        extracted: List[Dict[str, Any]] = []
        for table_index, raw_table in enumerate(page.extract_tables() or []):
            rows = []
            for row in raw_table or []:
                cleaned_row = [str(cell or "").strip() for cell in row]
                if any(cleaned_row):
                    rows.append(cleaned_row)
            if len(rows) < 2:
                continue
            extracted.append(
                {
                    "title": f"Table {page_number}.{table_index + 1}",
                    "headers": rows[0],
                    "rows": rows[1:],
                    "page": page_number,
                }
            )
        return extracted

    def _extract_images(self, page: Any, page_number: int) -> List[Dict[str, Any]]:
        images: List[Dict[str, Any]] = []
        for image_index, image in enumerate(page.images or []):
            images.append(
                {
                    "title": f"Image {page_number}.{image_index + 1}",
                    "page": page_number,
                    "width": image.get("width"),
                    "height": image.get("height"),
                    "x0": image.get("x0"),
                    "y0": image.get("y0"),
                }
            )
        return images

    def _assemble_structure(
        self,
        pages: List[Dict[str, Any]],
        heading_threshold: float,
        fallback_title: str,
    ) -> Dict[str, Any]:
        document_title = fallback_title.replace("_", " ").strip().title()
        modules: List[Dict[str, Any]] = []
        current_module: Optional[Dict[str, Any]] = None
        current_topic: Optional[Dict[str, Any]] = None

        def ensure_module(title: str = "Module 1") -> Dict[str, Any]:
            nonlocal current_module
            if current_module is None:
                current_module = {"title": title, "topics": [], "metadata": {}}
                modules.append(current_module)
            return current_module

        def ensure_topic(title: str = "Overview") -> Dict[str, Any]:
            nonlocal current_topic
            module = ensure_module()
            if current_topic is None:
                current_topic = {
                    "title": title,
                    "content_parts": [],
                    "tables": [],
                    "images": [],
                    "metadata": {},
                }
                module["topics"].append(current_topic)
            return current_topic

        def flush_topic():
            nonlocal current_topic
            if current_topic is None:
                return
            current_topic["content_text"] = self._normalize_text(" ".join(current_topic.pop("content_parts", [])))
            current_topic = None

        for page in pages:
            for line in page["lines"]:
                text = self._normalize_text(line["text"])
                if not text:
                    continue
                if self._is_heading(text, line.get("size", 12.0), heading_threshold):
                    if self._looks_like_module_heading(text):
                        flush_topic()
                        current_module = {"title": text, "topics": [], "metadata": {}}
                        modules.append(current_module)
                        current_topic = None
                    else:
                        flush_topic()
                        current_topic = {
                            "title": text,
                            "content_parts": [],
                            "tables": [],
                            "images": [],
                            "metadata": {},
                        }
                        ensure_module()["topics"].append(current_topic)
                else:
                    ensure_topic()["content_parts"].append(text)

            topic = ensure_topic()
            topic["tables"].extend(page["tables"])
            topic["images"].extend(page["images"])

        flush_topic()
        modules = [module for module in modules if module.get("topics")]
        if not modules:
            modules = [
                {
                    "title": "Module 1",
                    "topics": [
                        {
                            "title": "Overview",
                            "content_text": "",
                            "tables": [],
                            "images": [],
                            "metadata": {},
                        }
                    ],
                    "metadata": {},
                }
            ]

        for module in modules:
            for topic in module.get("topics", []):
                topic.setdefault("content_text", self._normalize_text(topic.get("content_text") or ""))
                topic.pop("content_parts", None)

        return {"title": document_title or "Unit", "modules": modules}

    def _is_heading(self, text: str, font_size: float, heading_threshold: float) -> bool:
        if len(text) < 4:
            return False
        if re.fullmatch(r"[\d\s.\-]+", text):
            return False
        if re.match(r"^(unit|module|chapter)\s+\d+", text, re.IGNORECASE):
            return True
        if (
            re.match(r"^\d+(\.\d+)*\s+[A-Z]", text)
            and len(text) <= 80
            and len(text.split()) <= 10
            and not text.endswith(".")
        ):
            return True
        if font_size < heading_threshold:
            return False
        if text.isupper() and len(text) <= 80 and len(text.split()) <= 8:
            return True
        if font_size >= heading_threshold and len(text) <= 120:
            return True
        if len(text.split()) <= 8 and text == text.title() and not text.endswith("."):
            return True
        return False

    def _looks_like_module_heading(self, text: str) -> bool:
        return bool(re.match(r"^(unit|module|chapter)\b", text, re.IGNORECASE))

    def _normalize_text(self, text: str) -> str:
        return re.sub(r"\s+", " ", text or "").strip()


class UnitEnrichmentService:
    def __init__(self):
        self.llm = get_llm_provider(feature="content_gen", provider="ollama")

    def enrich_topic(self, unit_id: str, topic: Dict[str, Any]) -> Dict[str, Any]:
        title = topic.get("title") or "Topic"
        content_text = (topic.get("content_text") or "").strip()
        source_tables = topic.get("source_tables") or []
        source_images = topic.get("source_images") or []
        signals = self._detect_signals(title, content_text)
        bullets = self._generate_bullets(title, content_text)
        quiz = self._generate_quiz(title, content_text)
        assets: List[Dict[str, Any]] = []

        if not source_tables and "comparison" in signals:
            table_payload = self._generate_table(title, content_text)
            if table_payload:
                assets.append(self._persist_table_asset(unit_id, topic["id"], table_payload))

        if not source_images and "process" in signals:
            steps = self._generate_process_steps(title, content_text)
            if steps:
                assets.append(self._persist_diagram_asset(unit_id, topic["id"], title, steps, "process"))
        elif not source_images and "concept" in signals:
            nodes = self._generate_concept_nodes(title, content_text)
            if nodes:
                assets.append(self._persist_diagram_asset(unit_id, topic["id"], title, nodes, "concept"))

        return {
            "summary_bullets": bullets,
            "quiz_questions": quiz,
            "detected_signals": signals,
            "assets": [asset for asset in assets if asset],
        }

    def _detect_signals(self, title: str, content_text: str) -> List[str]:
        haystack = f"{title} {content_text}".lower()
        signals: List[str] = []
        if any(keyword in haystack for keyword in ["process", "cycle", "workflow", "steps", "pipeline", "life cycle"]):
            signals.append("process")
        if any(keyword in haystack for keyword in ["compare", "comparison", "difference", "versus", "vs", "advantages", "disadvantages"]):
            signals.append("comparison")
        if any(keyword in haystack for keyword in ["components", "architecture", "types of", "classification", "framework", "system"]):
            signals.append("concept")
        return signals

    def _generate_bullets(self, topic_title: str, content_text: str) -> List[str]:
        prompt = (
            "Return JSON only. Provide 4 to 6 concise slide bullets for this topic.\n"
            f"Topic: {topic_title}\n"
            f"Content:\n{content_text[:5000]}\n"
            'Format: {"bullets": ["...", "..."]}'
        )
        payload = self._llm_json(prompt)
        bullets = payload.get("bullets") if isinstance(payload, dict) else None
        if isinstance(bullets, list) and bullets:
            return [str(item).strip() for item in bullets if str(item).strip()][:6]
        return self._fallback_bullets(content_text)

    def _generate_quiz(self, topic_title: str, content_text: str) -> List[Dict[str, Any]]:
        prompt = (
            "Return JSON only. Create 3 short quiz questions with answers.\n"
            f"Topic: {topic_title}\n"
            f"Content:\n{content_text[:5000]}\n"
            'Format: {"questions": [{"question":"...","answer":"..."}]}'
        )
        payload = self._llm_json(prompt)
        questions = payload.get("questions") if isinstance(payload, dict) else None
        if isinstance(questions, list) and questions:
            normalized = []
            for item in questions[:5]:
                if isinstance(item, dict) and item.get("question") and item.get("answer"):
                    normalized.append(
                        {"question": str(item["question"]).strip(), "answer": str(item["answer"]).strip()}
                    )
            if normalized:
                return normalized
        fallback = self._fallback_bullets(content_text)[:3]
        return [{"question": f"What is the key idea behind {topic_title}?", "answer": bullet} for bullet in fallback]

    def _generate_table(self, topic_title: str, content_text: str) -> Optional[Dict[str, Any]]:
        prompt = (
            "Return JSON only. Build a compact educational comparison table with 3 columns and up to 4 rows.\n"
            f"Topic: {topic_title}\n"
            f"Content:\n{content_text[:5000]}\n"
            'Format: {"title":"...","headers":["A","B","C"],"rows":[["...","...","..."]]}'
        )
        payload = self._llm_json(prompt)
        if isinstance(payload, dict) and isinstance(payload.get("headers"), list) and isinstance(payload.get("rows"), list):
            headers = [str(item).strip() for item in payload["headers"] if str(item).strip()]
            rows = [
                [str(cell).strip() for cell in row]
                for row in payload["rows"]
                if isinstance(row, list) and any(str(cell).strip() for cell in row)
            ]
            if headers and rows:
                return {"title": payload.get("title") or f"{topic_title} Comparison", "headers": headers, "rows": rows}
        bullets = self._fallback_bullets(content_text)[:3]
        if not bullets:
            return None
        return {
            "title": f"{topic_title} Comparison",
            "headers": ["Aspect", "Explanation", "Key Point"],
            "rows": [[f"Point {index + 1}", bullet, bullet.split(" ")[0]] for index, bullet in enumerate(bullets)],
        }

    def _generate_process_steps(self, topic_title: str, content_text: str) -> List[str]:
        prompt = (
            "Return JSON only. Extract 4 to 6 sequential process steps.\n"
            f"Topic: {topic_title}\n"
            f"Content:\n{content_text[:5000]}\n"
            'Format: {"steps":["step 1","step 2"]}'
        )
        payload = self._llm_json(prompt)
        steps = payload.get("steps") if isinstance(payload, dict) else None
        if isinstance(steps, list) and steps:
            return [str(item).strip() for item in steps if str(item).strip()][:6]
        return self._fallback_bullets(content_text)[:5]

    def _generate_concept_nodes(self, topic_title: str, content_text: str) -> List[str]:
        prompt = (
            "Return JSON only. Extract 4 to 6 concept labels closely related to the topic.\n"
            f"Topic: {topic_title}\n"
            f"Content:\n{content_text[:5000]}\n"
            'Format: {"nodes":["node 1","node 2"]}'
        )
        payload = self._llm_json(prompt)
        nodes = payload.get("nodes") if isinstance(payload, dict) else None
        if isinstance(nodes, list) and nodes:
            return [str(item).strip() for item in nodes if str(item).strip()][:6]
        return self._keyword_nodes(topic_title, content_text)

    def _fallback_bullets(self, content_text: str) -> List[str]:
        candidates = re.split(r"(?<=[.!?])\s+", content_text)
        bullets = [re.sub(r"\s+", " ", sentence).strip(" -") for sentence in candidates if len(sentence.strip()) > 20]
        return bullets[:5] or ["Topic content parsed and ready for review."]

    def _keyword_nodes(self, topic_title: str, content_text: str) -> List[str]:
        words = re.findall(r"[A-Za-z][A-Za-z-]{3,}", content_text.lower())
        common = [
            word.title()
            for word, _count in Counter(words).most_common(8)
            if word.lower() not in {"this", "that", "with", "from", "into", "their", "there", "topic"}
        ]
        return [topic_title] + common[:4]

    def _llm_json(self, prompt: str) -> Dict[str, Any]:
        try:
            response = self.llm.generate(
                prompt,
                system_prompt="You create strict JSON for educational content generation. Return JSON only.",
            )
            if not response or is_provider_error(response):
                return {}
            response = response.strip()
            if response.startswith("```"):
                response = re.sub(r"^```(?:json)?", "", response).strip()
                response = re.sub(r"```$", "", response).strip()
            return json.loads(response)
        except Exception:
            return {}

    def _persist_table_asset(self, unit_id: str, topic_id: str, table_payload: Dict[str, Any]) -> Dict[str, Any]:
        headers = table_payload.get("headers") or []
        rows = table_payload.get("rows") or []
        title = table_payload.get("title") or "Generated Table"

        html = self._table_to_html(title, headers, rows)
        csv_text = self._table_to_csv(headers, rows)
        base_name = self._slugify(title)

        html_url = storage_service.upload_bytes(
            html.encode("utf-8"),
            f"units/{unit_id}/assets/{topic_id}/{base_name}.html",
            content_type="text/html",
        )
        csv_url = storage_service.upload_bytes(
            csv_text.encode("utf-8"),
            f"units/{unit_id}/assets/{topic_id}/{base_name}.csv",
            content_type="text/csv",
        )
        return {
            "type": "table",
            "title": title,
            "file_url": html_url,
            "is_generated": True,
            "content_json": {"title": title, "headers": headers, "rows": rows},
            "metadata": {"csv_url": csv_url},
        }

    def _persist_diagram_asset(
        self,
        unit_id: str,
        topic_id: str,
        title: str,
        items: List[str],
        diagram_kind: str,
    ) -> Dict[str, Any]:
        png_bytes = self._build_diagram_png(title, items, diagram_kind)
        svg_text = self._build_diagram_svg(title, items, diagram_kind)
        base_name = self._slugify(title)

        png_url = storage_service.upload_bytes(
            png_bytes,
            f"units/{unit_id}/assets/{topic_id}/{base_name}.png",
            content_type="image/png",
        )
        svg_url = storage_service.upload_bytes(
            svg_text.encode("utf-8"),
            f"units/{unit_id}/assets/{topic_id}/{base_name}.svg",
            content_type="image/svg+xml",
        )
        return {
            "type": "diagram",
            "title": title,
            "file_url": png_url,
            "is_generated": True,
            "content_json": {"title": title, "items": items, "diagram_kind": diagram_kind},
            "metadata": {"svg_url": svg_url},
        }

    def _table_to_html(self, title: str, headers: List[str], rows: List[List[str]]) -> str:
        header_html = "".join(f"<th>{header}</th>" for header in headers)
        body_html = "".join(
            "<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>"
            for row in rows
        )
        return (
            "<html><head><meta charset='utf-8'>"
            "<style>body{font-family:Arial,sans-serif;padding:24px;}table{border-collapse:collapse;width:100%;}"
            "th,td{border:1px solid #d1d5db;padding:10px;text-align:left;}th{background:#ecfccb;}</style>"
            f"</head><body><h2>{title}</h2><table><thead><tr>{header_html}</tr></thead><tbody>{body_html}</tbody></table></body></html>"
        )

    def _table_to_csv(self, headers: List[str], rows: List[List[str]]) -> str:
        buffer = io.StringIO()
        writer = csv.writer(buffer)
        writer.writerow(headers)
        for row in rows:
            writer.writerow(row)
        return buffer.getvalue()

    def _build_diagram_png(self, title: str, items: List[str], diagram_kind: str) -> bytes:
        width, height = (1400, 900)
        image = Image.new("RGB", (width, height), "#f8fafc")
        draw = ImageDraw.Draw(image)
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()

        draw.rounded_rectangle((40, 30, width - 40, 110), radius=20, fill="#14532d")
        draw.text((70, 60), title[:90], fill="white", font=font_title)

        if diagram_kind == "process":
            top = 170
            box_height = 90
            for index, item in enumerate(items[:6]):
                left = 120 + (index % 2) * 620
                current_top = top + (index // 2) * 180
                draw.rounded_rectangle(
                    (left, current_top, left + 500, current_top + box_height),
                    radius=18,
                    fill="#dcfce7",
                    outline="#166534",
                    width=3,
                )
                draw.text((left + 20, current_top + 30), f"{index + 1}. {item[:55]}", fill="#052e16", font=font_body)
                if index < min(len(items[:6]) - 1, 5):
                    arrow_start = (left + 500, current_top + box_height // 2)
                    arrow_end = (left + 560, current_top + box_height // 2)
                    draw.line([arrow_start, arrow_end], fill="#166534", width=6)
                    draw.polygon(
                        [(arrow_end[0], arrow_end[1]), (arrow_end[0] - 20, arrow_end[1] - 10), (arrow_end[0] - 20, arrow_end[1] + 10)],
                        fill="#166534",
                    )
        else:
            center = (width // 2, height // 2)
            draw.ellipse((center[0] - 120, center[1] - 70, center[0] + 120, center[1] + 70), fill="#14532d")
            draw.text((center[0] - 80, center[1] - 10), title[:24], fill="white", font=font_body)
            positions = [(200, 180), (1120, 180), (200, 650), (1120, 650), (700, 170), (700, 690)]
            for index, item in enumerate(items[:6]):
                x, y = positions[index]
                draw.rounded_rectangle((x - 140, y - 45, x + 140, y + 45), radius=18, fill="#dcfce7", outline="#166534", width=3)
                draw.text((x - 110, y - 10), item[:28], fill="#052e16", font=font_body)
                draw.line([center, (x, y)], fill="#166534", width=4)

        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        return buffer.getvalue()

    def _build_diagram_svg(self, title: str, items: List[str], diagram_kind: str) -> str:
        if diagram_kind == "process":
            parts = [
                "<svg xmlns='http://www.w3.org/2000/svg' width='1400' height='900' viewBox='0 0 1400 900'>",
                "<rect x='40' y='30' width='1320' height='80' rx='20' fill='#14532d' />",
                f"<text x='70' y='78' font-size='28' fill='white'>{self._escape_xml(title)}</text>",
            ]
            top = 170
            for index, item in enumerate(items[:6]):
                x = 120 if index % 2 == 0 else 740
                y = top + (index // 2) * 180
                parts.append(f"<rect x='{x}' y='{y}' width='500' height='90' rx='18' fill='#dcfce7' stroke='#166534' stroke-width='3' />")
                parts.append(f"<text x='{x + 20}' y='{y + 52}' font-size='24' fill='#052e16'>{self._escape_xml(f'{index + 1}. {item[:55]}')}</text>")
            parts.append("</svg>")
            return "".join(parts)

        parts = [
            "<svg xmlns='http://www.w3.org/2000/svg' width='1400' height='900' viewBox='0 0 1400 900'>",
            "<rect x='0' y='0' width='1400' height='900' fill='#f8fafc' />",
            f"<ellipse cx='700' cy='450' rx='140' ry='85' fill='#14532d' />",
            f"<text x='620' y='458' font-size='28' fill='white'>{self._escape_xml(title[:24])}</text>",
        ]
        positions = [(200, 180), (1120, 180), (200, 650), (1120, 650), (700, 170), (700, 690)]
        for index, item in enumerate(items[:6]):
            x, y = positions[index]
            parts.append(f"<line x1='700' y1='450' x2='{x}' y2='{y}' stroke='#166534' stroke-width='4' />")
            parts.append(f"<rect x='{x - 140}' y='{y - 45}' width='280' height='90' rx='18' fill='#dcfce7' stroke='#166534' stroke-width='3' />")
            parts.append(f"<text x='{x - 110}' y='{y + 6}' font-size='22' fill='#052e16'>{self._escape_xml(item[:28])}</text>")
        parts.append("</svg>")
        return "".join(parts)

    def _slugify(self, value: str) -> str:
        slug = re.sub(r"[^\w\s-]", "", value).strip().lower()
        return re.sub(r"[-\s]+", "-", slug) or "asset"

    def _escape_xml(self, value: str) -> str:
        return (
            value.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
            .replace("'", "&apos;")
        )


class UnitPresentationService:
    def __init__(self):
        self.generator = PPTGenerator()

    def build_and_store_presentation(self, unit_detail: Dict[str, Any]) -> str:
        slides = []
        for module in unit_detail.get("modules") or []:
            for topic in module.get("topics") or []:
                slides.append(
                    {
                        "title": topic.get("title") or "Topic",
                        "bullets": (topic.get("summary_bullets") or [])[:5],
                    }
                )

        content_structure = {
            "title": unit_detail.get("title") or "Unit",
            "subtitle": f"{len(slides)} topic slides generated from uploaded PDF",
            "slides": slides,
        }
        filepath = self.generator.create_presentation(unit_detail.get("title") or "Unit", content_structure)

        presentation = Presentation(filepath)
        for module in unit_detail.get("modules") or []:
            for topic in module.get("topics") or []:
                for asset in topic.get("assets") or []:
                    if asset.get("type") == "diagram" and asset.get("file_url"):
                        self._append_diagram_slide(presentation, topic.get("title") or "Diagram", asset)
                    elif asset.get("type") == "table":
                        self._append_table_slide(presentation, topic.get("title") or "Table", asset)
        presentation.save(filepath)

        with open(filepath, "rb") as file_obj:
            ppt_bytes = file_obj.read()
        storage_url = storage_service.upload_bytes(
            ppt_bytes,
            f"units/{unit_detail['id']}/presentations/{self._slugify(unit_detail.get('title') or 'unit')}.pptx",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        try:
            os.remove(filepath)
        except OSError:
            pass
        return storage_url

    def _append_diagram_slide(self, presentation: Presentation, topic_title: str, asset: Dict[str, Any]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"{topic_title} Diagram"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_file:
            temp_path = temp_file.name
        try:
            storage_service.download_file(asset["file_url"], temp_path)
            slide.shapes.add_picture(temp_path, Inches(0.75), Inches(1.3), width=Inches(8.5))
        finally:
            try:
                os.remove(temp_path)
            except OSError:
                pass

    def _append_table_slide(self, presentation: Presentation, topic_title: str, asset: Dict[str, Any]) -> None:
        payload = asset.get("content_json") or {}
        headers = payload.get("headers") or []
        rows = payload.get("rows") or []
        if not headers:
            return
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"{topic_title} Table"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True

        shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.3), Inches(9), Inches(4.5))
        table = shape.table
        for column_index, header in enumerate(headers):
            table.cell(0, column_index).text = str(header)
        for row_index, row in enumerate(rows, start=1):
            for column_index, cell in enumerate(row[: len(headers)]):
                table.cell(row_index, column_index).text = str(cell)

    def _slugify(self, value: str) -> str:
        slug = re.sub(r"[^\w\s-]", "", value).strip().lower()
        return re.sub(r"[-\s]+", "-", slug) or "presentation"


unit_pdf_parser = UnitPDFParser()
unit_enrichment_service = UnitEnrichmentService()
unit_presentation_service = UnitPresentationService()
