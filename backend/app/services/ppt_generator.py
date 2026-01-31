"""
PowerPoint Generation Service
Generates professional PPTX presentations using python-pptx library.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime
from typing import List, Dict, Any
import re


class PPTGenerator:
    """Service for generating PowerPoint presentations."""
    
    # Lumina Brand Colors
    BRAND_PRIMARY = RGBColor(34, 197, 94)  # Green #22C55E
    BRAND_SECONDARY = RGBColor(168, 85, 247)  # Purple #A855F7
    DARK_BG = RGBColor(15, 17, 21)  # Dark background #0F1115
    LIGHT_TEXT = RGBColor(243, 244, 246)  # Light gray #F3F4F6
    GRAY_TEXT = RGBColor(156, 163, 175)  # Gray #9CA3AF
    
    def __init__(self, output_dir: str = "static/presentations"):
        """Initialize the PPT generator.
        
        Args:
            output_dir: Directory to save generated presentations
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def create_presentation(self, topic: str, content_structure: Dict[str, Any]) -> str:
        """Create a PowerPoint presentation from structured content.
        
        Args:
            topic: Main topic/title of the presentation
            content_structure: Dictionary with title, subtitle, and slides
            
        Returns:
            Path to the generated PPTX file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self._add_title_slide(
            prs, 
            content_structure.get("title", topic),
            content_structure.get("subtitle", f"A comprehensive overview of {topic}")
        )
        
        # Add content slides
        for slide_data in content_structure.get("slides", []):
            self._add_content_slide(
                prs,
                slide_data.get("title", ""),
                slide_data.get("bullets", [])
            )
        
        # Save presentation
        filename = self._generate_filename(topic)
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        
        return filepath
    
    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str):
        """Add a title slide with Lumina branding.
        
        Args:
            prs: Presentation object
            title: Main title text
            subtitle: Subtitle text
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.DARK_BG
        
        # Add gradient accent (top-right corner)
        left = Inches(7)
        top = Inches(0)
        width = Inches(3)
        height = Inches(3)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.BRAND_PRIMARY
        shape.fill.transparency = 0.9
        shape.line.fill.background()
        
        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.LIGHT_TEXT
        
        # Add subtitle
        left = Inches(1)
        top = Inches(4.2)
        width = Inches(8)
        height = Inches(1)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = self.GRAY_TEXT
        
        # Add Lumina branding
        left = Inches(4)
        top = Inches(6.5)
        width = Inches(2)
        height = Inches(0.5)
        brand_box = slide.shapes.add_textbox(left, top, width, height)
        brand_frame = brand_box.text_frame
        brand_frame.text = "Powered by Lumina AI"
        brand_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        brand_frame.paragraphs[0].font.size = Pt(12)
        brand_frame.paragraphs[0].font.color.rgb = self.BRAND_PRIMARY
    
    def _add_content_slide(self, prs: Presentation, title: str, bullets: List[str]):
        """Add a content slide with title and bullet points.
        
        Args:
            prs: Presentation object
            title: Slide title
            bullets: List of bullet point texts
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.DARK_BG
        
        # Add accent bar (left side)
        left = Inches(0)
        top = Inches(0)
        width = Inches(0.1)
        height = Inches(7.5)
        accent = slide.shapes.add_shape(1, left, top, width, height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.BRAND_PRIMARY
        accent.line.fill.background()
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.LIGHT_TEXT
        
        # Add bullet points
        left = Inches(1)
        top = Inches(2)
        width = Inches(8.5)
        height = Inches(5)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet_text
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = self.GRAY_TEXT
            p.space_before = Pt(12)
            
            # Add bullet character
            p.text = f"• {bullet_text}"
    
    def _generate_filename(self, topic: str) -> str:
        """Generate a safe filename from topic.
        
        Args:
            topic: Presentation topic
            
        Returns:
            Safe filename string
        """
        # Remove special characters and replace spaces with underscores
        safe_topic = re.sub(r'[^\w\s-]', '', topic)
        safe_topic = re.sub(r'[-\s]+', '_', safe_topic)
        safe_topic = safe_topic.lower()[:50]  # Limit length
        
        # Add timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        return f"{safe_topic}_{timestamp}.pptx"
    
    def get_file_size(self, filepath: str) -> str:
        """Get human-readable file size.
        
        Args:
            filepath: Path to file
            
        Returns:
            File size string (e.g., "2.3 MB")
        """
        size_bytes = os.path.getsize(filepath)
        
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.1f} KB"
        else:
            return f"{size_bytes / (1024 * 1024):.1f} MB"
